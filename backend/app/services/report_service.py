from app.pdf_builders.base_builder import get_builder
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io

def generate_pdf(analysis: dict, report_type: str = "ventas", profile: dict = None) -> bytes:
    builder = get_builder(report_type, analysis, profile)
    return builder.build_pdf()

def generate_pptx(analysis: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    TEAL_RGB = RGBColor(0x1D, 0x9E, 0x75)
    DARK_RGB = RGBColor(0x1a, 0x1a, 0x2e)
    WHITE_RGB = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY_RGB = RGBColor(0x6b, 0x72, 0x80)

    def add_slide(layout_idx=6):
        return prs.slides.add_slide(prs.slide_layouts[layout_idx])

    def add_textbox(slide, text, left, top, width, height,
                    size=18, bold=False, color=DARK_RGB, wrap=True):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text = text
        p.runs[0].font.size = Pt(size)
        p.runs[0].font.bold = bold
        p.runs[0].font.color.rgb = color
        return txBox

    slide1 = add_slide()
    bg = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
    bg.line.fill.background()
    accent = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL_RGB
    accent.line.fill.background()
    add_textbox(slide1, analysis['titulo'], 0.5, 2.5, 12, 1.2, size=32, bold=True, color=DARK_RGB)
    add_textbox(slide1, 'Reporte generado por Reporti', 0.5, 3.9, 8, 0.5, size=14, color=GRAY_RGB)

    slide2 = add_slide()
    add_textbox(slide2, 'Resumen ejecutivo', 0.4, 0.3, 12, 0.6, size=20, bold=True, color=TEAL_RGB)
    add_textbox(slide2, analysis['resumen_ejecutivo'], 0.4, 1.1, 12.5, 5, size=16, color=DARK_RGB)

    slide3 = add_slide()
    add_textbox(slide3, 'KPIs principales', 0.4, 0.3, 12, 0.6, size=20, bold=True, color=TEAL_RGB)
    cols = min(len(analysis['kpis']), 4)
    col_w = 12.5 / cols
    for i, kpi in enumerate(analysis['kpis'][:4]):
        x = 0.4 + i * col_w
        box = slide3.shapes.add_shape(1, Inches(x), Inches(1.2), Inches(col_w - 0.2), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xE1, 0xF5, 0xEE)
        box.line.color.rgb = TEAL_RGB
        add_textbox(slide3, kpi['valor'], x+0.1, 1.5, col_w-0.3, 0.9, size=22, bold=True, color=TEAL_RGB)
        add_textbox(slide3, kpi['nombre'], x+0.1, 2.4, col_w-0.3, 0.5, size=11, bold=True, color=DARK_RGB)
        add_textbox(slide3, kpi['descripcion'], x+0.1, 2.9, col_w-0.3, 0.7, size=9, color=GRAY_RGB)

    slide4 = add_slide()
    add_textbox(slide4, 'Recomendaciones', 0.4, 0.3, 12, 0.6, size=20, bold=True, color=TEAL_RGB)
    for i, rec in enumerate(analysis['recomendaciones'][:4]):
        y = 1.2 + i * 1.4
        dot = slide4.shapes.add_shape(9, Inches(0.4), Inches(y+0.15), Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = TEAL_RGB
        dot.line.fill.background()
        add_textbox(slide4, rec['titulo'], 0.75, y, 12, 0.4, size=13, bold=True, color=DARK_RGB)
        add_textbox(slide4, rec['descripcion'], 0.75, y+0.4, 12, 0.8, size=11, color=GRAY_RGB)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()